# Copyright (C) 2024 Intel Corporation
# SPDX-License-Identifier: Apache-2.0
import asyncio
import base64
import errno
import functools
import io
import json
import multiprocessing
import os
import re
import shutil
import signal
import subprocess
import tempfile
import timeit
import unicodedata
import urllib.parse
from pathlib import Path
from typing import Dict, List, Union
from urllib.parse import urlparse, urlunparse

import docx
import docx2txt

import fitz
import numpy as np
import pandas as pd
import pptx
import requests
import yaml
from bs4 import BeautifulSoup
from docx.oxml.ns import qn
from docx.table import Table, _Cell
from docx.text.paragraph import Paragraph
from langchain import LLMChain, PromptTemplate
# from langchain_community.document_loaders import (
#     UnstructuredHTMLLoader,
#     UnstructuredImageLoader,
#     UnstructuredMarkdownLoader,
#     UnstructuredXMLLoader,
# )
from langchain_community.llms import HuggingFaceEndpoint
from PIL import Image
from openai import OpenAI
from sklearn.metrics.pairwise import cosine_similarity

from comps import CustomLogger
from comps.dataprep.config import get_dataprep_llm_config, get_llm_extra_body
from comps.dataprep.prompt.prompt_manager import PromptRegistry, PromptKey, Lang

logger = CustomLogger("prepare_doc_util", os.getenv("LOG_LEVEL", "INFO"))

class TimeoutError(Exception):
    pass


def timeout(seconds=10, error_message=os.strerror(errno.ETIME)):
    def decorator(func):
        def _handle_timeout(signum, frame):
            raise TimeoutError(error_message)

        @functools.wraps(func)
        def wrapper(*args, **kwargs):
            signal.signal(signal.SIGALRM, _handle_timeout)
            signal.alarm(seconds)
            try:
                result = func(*args, **kwargs)
            finally:
                signal.alarm(0)
            return result

        return wrapper

    return decorator


class Timer:
    level = 0
    viewer = None

    def __init__(self, name):
        self.name = name
        if Timer.viewer:
            Timer.viewer.display(f"{name} started ...")
        else:
            print(f"{name} started ...")

    def __enter__(self):
        self.start = timeit.default_timer()
        Timer.level += 1

    def __exit__(self, *a, **kw):
        Timer.level -= 1
        if Timer.viewer:
            Timer.viewer.display(f'{"  " * Timer.level}{self.name} took {timeit.default_timer() - self.start} sec')
        else:
            print(f'{"  " * Timer.level}{self.name} took {timeit.default_timer() - self.start} sec')


def get_separators():
    separators = [
        "\n\n",
        "\n",
        " ",
        ".",
        ",",
        "\u200b",  # Zero-width space
        "\uff0c",  # Fullwidth comma
        "\u3001",  # Ideographic comma
        "\uff0e",  # Fullwidth full stop
        "\u3002",  # Ideographic full stop
        "",
    ]
    return separators


# def load_pdf(pdf_path):
#     """Load the pdf file."""
#     doc = fitz.open(pdf_path)
#     reader = easyocr.Reader(["en"], gpu=False)
#     result = ""
#     for i in range(doc.page_count):
#         page = doc.load_page(i)
#         pagetext = page.get_text().strip()
#         if pagetext:
#             if pagetext.endswith("!") or pagetext.endswith("?") or pagetext.endswith("."):
#                 result = result + pagetext
#             else:
#                 result = result + pagetext + "."
#         if len(doc.get_page_images(i)) > 0:
#             for img in doc.get_page_images(i):
#                 if img:
#                     pageimg = ""
#                     xref = img[0]
#                     img_data = doc.extract_image(xref)
#                     img_bytes = img_data["image"]
#                     pil_image = Image.open(io.BytesIO(img_bytes))
#                     img = np.array(pil_image)
#                     img_result = reader.readtext(img, paragraph=True, detail=0)
#                     pageimg = pageimg + ", ".join(img_result).strip()
#                     if pageimg.endswith("!") or pageimg.endswith("?") or pageimg.endswith("."):
#                         pass
#                     else:
#                         pageimg = pageimg + "."
#                 result = result + pageimg
#     return result


# def load_html(html_path):
#     """Load the html file."""
#     data_html = UnstructuredHTMLLoader(html_path).load()
#     content = ""
#     for ins in data_html:
#         content += ins.page_content
#     return content


def load_txt(txt_path):
    """Load txt file."""
    with open(txt_path, "r") as file:
        text = file.read()
    return text


# def load_doc(doc_path):
#     """Load doc file."""
#     print("Converting doc file to docx file...")
#     docx_path = doc_path + "x"
#     subprocess.run(
#         [
#             "libreoffice",
#             "--headless",
#             "--invisible",
#             "--convert-to",
#             "docx",
#             "--outdir",
#             os.path.dirname(docx_path),
#             doc_path,
#         ],
#         check=True,
#     )
#     print("Converted doc file to docx file.")
#     text = load_docx(docx_path)
#     os.remove(docx_path)
#     return text


# def load_docx(docx_path):
#     """Load docx file."""
#     doc = docx.Document(docx_path)
#     text = ""
#     # Save all 'rId:filenames' relationships in an dictionary and save the images if any.
#     rid2img = {}
#     for r in doc.part.rels.values():
#         if isinstance(r._target, docx.parts.image.ImagePart):
#             rid2img[r.rId] = os.path.basename(r._target.partname)
#     if rid2img:
#         save_path = tempfile.mkdtemp()
#         docx2txt.process(docx_path, save_path)
#     for paragraph in doc.paragraphs:
#         if hasattr(paragraph, "text"):
#             text += paragraph.text + "\n"
#         if "graphicData" in paragraph._p.xml:
#             for rid in rid2img:
#                 if rid in paragraph._p.xml:
#                     img_path = os.path.join(save_path, rid2img[rid])
#                     img_text = load_image(img_path)
#                     if img_text:
#                         text += img_text + "\n"
#     if rid2img:
#         shutil.rmtree(save_path)
#     return text


# def load_ppt(ppt_path):
#     """Load ppt file."""
#     print("Converting ppt file to pptx file...")
#     pptx_path = ppt_path + "x"
#     subprocess.run(
#         [
#             "libreoffice",
#             "--headless",
#             "--invisible",
#             "--convert-to",
#             "docx",
#             "--outdir",
#             os.path.dirname(pptx_path),
#             ppt_path,
#         ],
#         check=True,
#     )
#     print("Converted ppt file to pptx file.")
#     text = load_pptx(pptx_path)
#     os.remove(pptx_path)
#     return text


# def load_pptx(pptx_path):
#     """Load pptx file."""
#     text = ""
#     prs = pptx.Presentation(pptx_path)
#     for slide in prs.slides:
#         for shape in sorted(slide.shapes, key=lambda shape: (shape.top, shape.left)):
#             if shape.has_text_frame:
#                 if shape.text:
#                     text += shape.text + "\n"
#             if shape.has_table:
#                 table_contents = "\n".join(
#                     [
#                         "\t".join([(cell.text if hasattr(cell, "text") else "") for cell in row.cells])
#                         for row in shape.table.rows
#                         if hasattr(row, "cells")
#                     ]
#                 )
#                 if table_contents:
#                     text += table_contents + "\n"
#             if hasattr(shape, "image") and hasattr(shape.image, "blob"):
#                 img_path = f"./{shape.image.filename}"
#                 with open(img_path, "wb") as f:
#                     f.write(shape.image.blob)
#                 img_text = load_image(img_path)
#                 if img_text:
#                     text += img_text + "\n"
#                 os.remove(img_path)
#     return text


# def load_md(md_path):
#     """Load md file."""
#     loader = UnstructuredMarkdownLoader(md_path)
#     text = loader.load()[0].page_content
#     return text
#
#
# def load_xml(xml_path):
#     """Load xml file."""
#     loader = UnstructuredXMLLoader(xml_path)
#     text = loader.load()[0].page_content
#     return text


def load_json(json_path):
    """Load and process json file."""
    with open(json_path, "r") as file:
        data = json.load(file)
    content_list = [json.dumps(item) for item in data]
    return content_list


def load_jsonl(jsonl_path):
    """Load and process jsonl file."""
    content_list = []
    with open(jsonl_path, "r") as file:
        for line in file:
            json_obj = json.loads(line)
            content_list.append(json_obj)
    return content_list


def load_yaml(yaml_path):
    """Load and process yaml file."""
    with open(yaml_path, "r") as file:
        data = yaml.safe_load(file)
    return yaml.dump(data)


def load_xlsx(input_path):
    """Load and process xlsx file."""
    df = pd.read_excel(input_path)
    content_list = df.apply(lambda row: ", ".join(row.astype(str)), axis=1).tolist()
    return content_list


def load_csv(input_path):
    """Load the csv file."""
    df = pd.read_csv(input_path)
    content_list = df.apply(lambda row: ", ".join(row.astype(str)), axis=1).tolist()
    return content_list


# def load_image(image_path):
#     """Load the image file."""
#     if os.getenv("SUMMARIZE_IMAGE_VIA_LVM", None) == "1":
#         query = "Please summarize this image."
#         image_b64_str = base64.b64encode(open(image_path, "rb").read()).decode()
#         response = requests.post(
#             "http://localhost:9399/v1/lvm",
#             data=json.dumps({"image": image_b64_str, "prompt": query}),
#             headers={"Content-Type": "application/json"},
#             proxies={"http": None},
#         )
#         return response.json()["text"].strip()
#     loader = UnstructuredImageLoader(image_path)
#     text = loader.load()[0].page_content
#     return text.strip()


# def load_svg(svg_path):
#     """Load the svg file."""
#     png_path = svg_path.replace(".svg", ".png")
#     cairosvg.svg2png(url=svg_path, write_to=png_path)
#     text = load_image(png_path)
#     os.remove(png_path)
#     return text


def document_loader(doc_path):
    # if doc_path.endswith(".pdf"):
    #     return load_pdf(doc_path)
    # if doc_path.endswith(".html"):
    #     return load_html(doc_path)
    if doc_path.endswith(".txt"):
        return load_txt(doc_path)
    # elif doc_path.endswith(".doc"):
    #     return load_doc(doc_path)
    # elif doc_path.endswith(".docx"):
    #     return load_docx(doc_path)
    # elif doc_path.endswith(".ppt"):
    #     return load_ppt(doc_path)
    # elif doc_path.endswith(".pptx"):
    #     return load_pptx(doc_path)
    # elif doc_path.endswith(".md"):
    #     return load_md(doc_path)
    # elif doc_path.endswith(".xml"):
    #     return load_xml(doc_path)
    elif doc_path.endswith(".json"):
        return load_json(doc_path)
    elif doc_path.endswith(".jsonl"):
        return load_jsonl(doc_path)
    elif doc_path.endswith(".yaml"):
        return load_yaml(doc_path)
    elif doc_path.endswith(".xlsx") or doc_path.endswith(".xls"):
        return load_xlsx(doc_path)
    elif doc_path.endswith(".csv"):
        return load_csv(doc_path)
    # elif (
    #     doc_path.endswith(".tiff")
    #     or doc_path.endswith(".jpg")
    #     or doc_path.endswith(".jpeg")
    #     or doc_path.endswith(".png")
    # ):
    #     return load_image(doc_path)
    # elif doc_path.endswith(".svg"):
    #     return load_svg(doc_path)
    else:
        raise NotImplementedError(
            "Current only support pdf, html, txt, doc, docx, pptx, ppt, md, xml"
            + ", json, jsonl, yaml, xlsx, xls, csv, tiff and svg format."
        )


class Crawler:
    def __init__(self, pool=None):
        if pool:
            assert isinstance(pool, (str, list, tuple)), "url pool should be str, list or tuple"
        self.pool = pool
        self.headers = {
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng, \
            */*;q=0.8,application/signed-exchange;v=b3;q=0.7",
            "Accept-Encoding": "gzip, deflate, br",
            "Accept-Language": "en-US,en;q=0.9,zh-CN;q=0.8,zh;q=0.7",
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, \
            like Gecko) Chrome/113.0.0.0 Safari/537.36",
        }
        self.fetched_pool = set()

    def get_sublinks(self, soup):
        sublinks = []
        for links in soup.find_all("a"):
            sublinks.append(str(links.get("href")))
        return sublinks

    def get_hyperlink(self, soup, base_url):
        sublinks = []
        for links in soup.find_all("a"):
            link = str(links.get("href"))
            if link.startswith("#") or link is None or link == "None":
                continue
            suffix = link.split("/")[-1]
            if "." in suffix and suffix.split(".")[-1] not in ["html", "htmld"]:
                continue
            link_parse = urlparse(link)
            base_url_parse = urlparse(base_url)
            if link_parse.path == "":
                continue
            if link_parse.netloc != "":
                # keep crawler works in the same domain
                if link_parse.netloc != base_url_parse.netloc:
                    continue
                sublinks.append(link)
            else:
                sublinks.append(
                    urlunparse(
                        (
                            base_url_parse.scheme,
                            base_url_parse.netloc,
                            link_parse.path,
                            link_parse.params,
                            link_parse.query,
                            link_parse.fragment,
                        )
                    )
                )
        return sublinks

    def fetch(self, url, headers=None, max_times=5):
        if not headers:
            headers = self.headers
        while max_times:
            parsed_url = urlparse(url)
            if not parsed_url.scheme:
                url = "http://" + url

            logger.debug("start fetch %s..." % url)
            try:
                response = requests.get(url, headers=headers, verify=True)
                if response.status_code != 200:
                    print("fail to fetch %s, response status code: %s", url, response.status_code)
                else:
                    # Extract charset from the Content-Type header
                    content_type = response.headers.get("Content-Type", "").lower()
                    if "charset=" in content_type:
                        # Extract charset value from the content-type header
                        charset = content_type.split("charset=")[-1].strip()
                        response.encoding = charset
                        logger.debug(f"Charset detected and set: {response.encoding}")
                    else:
                        import re

                        # Extract charset from the response HTML content
                        charset_from_meta = None
                        # Check for <meta charset="...">
                        match = re.search(r'<meta\s+charset=["\']?([^"\'>]+)["\']?', response.text, re.IGNORECASE)
                        if match:
                            charset_from_meta = match.group(1)
                        # Check for <meta http-equiv="Content-Type" content="...; charset=...">
                        if not charset_from_meta:
                            match = re.search(
                                r'<meta\s+http-equiv=["\']?content-type["\']?\s+content=["\']?[^"\']*charset=([^"\'>]+)["\']?',
                                response.text,
                                re.IGNORECASE,
                            )
                            if match:
                                charset_from_meta = match.group(1)
                        if charset_from_meta:
                            response.encoding = charset_from_meta
                            logger.debug(f"Charset detected and set from meta tag: {response.encoding}")
                        else:
                            # Fallback to default encoding
                            response.encoding = "utf-8"
                            logger.debug("Charset not specified, using default utf-8")
                    return response
            except Exception as e:
                print("fail to fetch %s, caused by %s", url, e)
                raise Exception(e)
            max_times -= 1
        return None

    def process_work(self, sub_url, work):
        response = self.fetch(sub_url)
        if response is None:
            return []
        self.fetched_pool.add(sub_url)
        soup = self.parse(response.text)
        base_url = self.get_base_url(sub_url)
        sublinks = self.get_hyperlink(soup, base_url)
        if work:
            work(sub_url, soup)
        return sublinks

    def crawl(self, pool, work=None, max_depth=10, workers=10):
        url_pool = set()
        for url in pool:
            base_url = self.get_base_url(url)
            response = self.fetch(url)
            soup = self.parse(response.text)
            sublinks = self.get_hyperlink(soup, base_url)
            self.fetched_pool.add(url)
            url_pool.update(sublinks)
            depth = 0
            while len(url_pool) > 0 and depth < max_depth:
                print("current depth %s...", depth)
                mp = multiprocessing.Pool(processes=workers)
                results = []
                for sub_url in url_pool:
                    if sub_url not in self.fetched_pool:
                        results.append(mp.apply_async(self.process_work, (sub_url, work)))
                mp.close()
                mp.join()
                url_pool = set()
                for result in results:
                    sublinks = result.get()
                    url_pool.update(sublinks)
                depth += 1

    def parse(self, html_doc):
        soup = BeautifulSoup(html_doc, "lxml")
        return soup

    def download(self, url, file_name):
        print("download %s into %s...", url, file_name)
        try:
            r = requests.get(url, stream=True, headers=self.headers, verify=True)
            f = open(file_name, "wb")
            for chunk in r.iter_content(chunk_size=512):
                if chunk:
                    f.write(chunk)
        except Exception as e:
            print("fail to download %s, caused by %s", url, e)

    def get_base_url(self, url):
        result = urlparse(url)
        return urlunparse((result.scheme, result.netloc, "", "", "", ""))

    def clean_text(self, text):
        text = text.strip().replace("\r", "\n")
        text = re.sub(" +", " ", text)
        text = re.sub("\n+", "\n", text)
        text = text.split("\n")
        return "\n".join([i for i in text if i and i != " "])


def uni_pro(text):
    """Check if the character is ASCII or falls in the category of non-spacing marks."""
    normalized_text = unicodedata.normalize("NFKD", text)
    filtered_text = ""
    for char in normalized_text:
        if ord(char) < 128 or unicodedata.category(char) == "Mn":
            filtered_text += char
    return filtered_text


def load_html_data(url):
    crawler = Crawler()
    res = crawler.fetch(url)
    if res is None:
        return None
    soup = crawler.parse(res.text)
    all_text = crawler.clean_text(soup.select_one("body").text)
    main_content = ""
    for element_name in ["main", "container"]:
        main_block = None
        if soup.select(f".{element_name}"):
            main_block = soup.select(f".{element_name}")
        elif soup.select(f"#{element_name}"):
            main_block = soup.select(f"#{element_name}")
        if main_block:
            for element in main_block:
                text = crawler.clean_text(element.text)
                if text not in main_content:
                    main_content += f"\n{text}"
            main_content = crawler.clean_text(main_content)
    main_content = all_text if main_content == "" else main_content
    main_content = main_content.replace("\n", "")
    main_content = main_content.replace("\n\n", "")
    main_content = re.sub(r"\s+", " ", main_content)
    logger.debug("main_content=[%s]" % main_content)

    return main_content


def parse_html(input):
    """Parse the uploaded file."""
    chucks = []
    for link in input:
        if re.match(r"^https?:/{2}\w.+$", link):
            content = load_html_data(link)
            if content is None:
                continue
            chuck = [[content.strip(), link]]
            chucks += chuck
        else:
            print("The given link/str {} cannot be parsed.".format(link))

    return chucks


# def get_tables_result(pdf_path, table_strategy):
#     """Extract tables information from pdf file."""
#     if table_strategy == "fast":
#         return None
#
#     from unstructured.documents.elements import FigureCaption
#     from unstructured.partition.pdf import partition_pdf
#
#     tables_result = []
#     raw_pdf_elements = partition_pdf(
#         filename=pdf_path,
#         infer_table_structure=True,
#     )
#     tables = [el for el in raw_pdf_elements if el.category == "Table"]
#     for table in tables:
#         table_coords = table.metadata.coordinates.points
#         content = table.metadata.text_as_html
#         table_page_number = table.metadata.page_number
#         min_distance = float("inf")
#         table_summary = None
#         if table_strategy == "hq":
#             for element in raw_pdf_elements:
#                 if isinstance(element, FigureCaption) or element.text.startswith("Tab"):
#                     caption_page_number = element.metadata.page_number
#                     caption_coords = element.metadata.coordinates.points
#                     related, y_distance = get_relation(
#                         table_coords, caption_coords, table_page_number, caption_page_number
#                     )
#                     if related:
#                         if y_distance < min_distance:
#                             min_distance = y_distance
#                             table_summary = element.text
#             if table_summary is None:
#                 parent_id = table.metadata.parent_id
#                 for element in raw_pdf_elements:
#                     if element.id == parent_id:
#                         table_summary = element.text
#                         break
#         elif table_strategy == "llm":
#             table_summary = llm_generate(content)
#             table_summary = table_summary.lstrip("\n ")
#         elif table_strategy is None:
#             table_summary = None
#         if table_summary is None:
#             text = f"[Table: {content}]"
#         else:
#             text = f"|Table: [Summary: {table_summary}], [Content: {content}]|"
#         tables_result.append(text)
#     return tables_result


def llm_generate(content):
    llm_endpoint = os.getenv("TGI_LLM_ENDPOINT", "http://localhost:8080")
    llm = HuggingFaceEndpoint(
        endpoint_url=llm_endpoint,
        max_new_tokens=1000,
        top_k=40,
        top_p=0.9,
        temperature=0.8,
        streaming=False,
        num_beams=2,
        num_return_sequences=2,
        use_cache=True,
        timeout=600,
    )

    table_summary_template = """
    Task: Your task is to give a concise summary of the table. \
    The summary should cover the overall table structure and all detailed information of the table. \
    The table will be given in html format. Summarize the table below.
    ---
    ### Table:
    {table_content}
    ---
    ### Generated Summary:
    """

    prompt = PromptTemplate(template=table_summary_template, input_variables=["table_content"])

    llm_chain = LLMChain(prompt=prompt, llm=llm)

    response = llm_chain.invoke(content)
    response = response["text"]
    print("response", response)
    return response


def get_relation(table_coords, caption_coords, table_page_number, caption_page_number, threshold=100):
    """Get the relation of a pair of table and caption."""
    same_page = table_page_number == caption_page_number
    x_overlap = (min(table_coords[2][0], caption_coords[2][0]) - max(table_coords[0][0], caption_coords[0][0])) > 0
    if table_coords[0][1] - caption_coords[1][1] >= 0:
        y_distance = table_coords[0][1] - caption_coords[1][1]
    elif caption_coords[0][1] - table_coords[1][1] >= 0:
        y_distance = caption_coords[0][1] - table_coords[1][1]
    else:
        y_distance = 0
    y_close = y_distance < threshold
    return same_page and x_overlap and y_close, y_distance


def create_upload_folder(upload_path):
    if not os.path.exists(upload_path):
        Path(upload_path).mkdir(parents=True, exist_ok=True)


def encode_filename(filename):
    return urllib.parse.quote(filename, safe="")


def decode_filename(encoded_filename):
    return urllib.parse.unquote(encoded_filename)


async def save_content_to_local_disk(save_path: str, content):
    save_path = Path(save_path)
    try:
        if save_path.exists():
            return {"status": 400, "message": f"File {save_path} already exists."}

        if isinstance(content, str):
            with open(save_path, "w", encoding="utf-8") as file:
                file.write(content)
        else:
            with save_path.open("wb") as fout:
                content = await content.read()
                fout.write(content)

        return {"status": "success", "message": f"File {save_path} saved successfully."}
    except Exception as e:
        error_message = f"Write file {save_path} failed. Exception: {e}"
        print(error_message)
        return {"status": 500, "message": error_message}


def get_file_structure(root_path: str, parent_path: str = "") -> List[Dict[str, Union[str, List]]]:
    result = []
    for path in os.listdir(root_path):
        complete_path = parent_path + "/" + path if parent_path else path
        file_path = root_path + "/" + path
        p = Path(file_path)
        # append file into result
        if p.is_file():
            file_dict = {
                "name": decode_filename(path),
                "id": decode_filename(complete_path),
                "type": "File",
                "parent": "",
            }
            result.append(file_dict)
        else:
            # append folder and inner files/folders into result using recursive function
            folder_dict = {
                "name": decode_filename(path),
                "id": decode_filename(complete_path),
                "type": "Directory",
                "children": get_file_structure(file_path, complete_path),
                "parent": "",
            }
            result.append(folder_dict)

    return result


def format_search_results(response, file_list: list):
    for i in range(1, len(response), 2):
        file_name = response[i].decode()[5:]
        file_dict = {
            "name": decode_filename(file_name),
            "id": decode_filename(file_name),
            "type": "File",
            "parent": "",
        }
        file_list.append(file_dict)
    return file_list


def remove_folder_with_ignore(folder_path: str, except_patterns: List = []):
    """Remove the specific folder, and ignore some files/folders.

    :param folder_path: file path to delete
    :param except_patterns: files/folder name to ignore
    """
    print(f"except patterns: {except_patterns}")
    for root, dirs, files in os.walk(folder_path, topdown=False):
        for name in files:
            # delete files except ones that match patterns
            file_path = os.path.join(root, name)
            if except_patterns != [] and any(pattern in file_path for pattern in except_patterns):
                continue
            os.remove(file_path)

        # delete empty folder
        for name in dirs:
            dir_path = os.path.join(root, name)
            # delete folders except ones that match patterns
            if except_patterns != [] and any(pattern in dir_path for pattern in except_patterns):
                continue
            if not os.listdir(dir_path):
                os.rmdir(dir_path)

def normalize(text):
    return str(text).strip().replace("　", "").replace(" ", "").replace("\n","")


def iter_block_items(parent):
    """
    按顺序遍历文档或表格单元格的段落和表格，保持原文结构顺序。
    支持：
    - docx.Document
    - _Cell（单元格内嵌表格）
    返回生成器，yield Paragraph 或 Table 对象
    """
    if isinstance(parent, docx.document.Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError(f"Unsupported parent type: {type(parent)}")

    for child in parent_elm.iterchildren():
        # 段落
        if child.tag == qn("w:p"):
            yield Paragraph(child, parent)
        # 表格
        elif child.tag == qn("w:tbl"):
            yield Table(child, parent)
        # 忽略其他类型，如注释、图片等，可以根据需要扩展


def replace_html_tables_with_placeholders(content: str, html_df_map: dict,embeddings) -> str:
    """
    将每个 <html> 区块中的表格替换为占位符，使用语义匹配找到对应的真实表格内容。
    采用滑动指针方式，按顺序匹配，不回退。
    """

    # 将 html_df_map 转换为列表，方便滑动指针操作
    html_list = list(html_df_map.keys())
    df_list = list(html_df_map.values())
    html_pointer = 0  # 滑动指针，只向前不回退

    table_counter = 1
    similarity_threshold = 0.75

    def html_replacer(match):
        nonlocal table_counter, html_pointer

        html_block = match.group()
        soup = BeautifulSoup(html_block, "html.parser")
        placeholders = []

        for table in soup.find_all("table"):
            table_html = str(table)
            best_match_content = None

            # 从当前指针位置开始向后查找匹配
            for i in range(html_pointer, len(html_list)):
                candidate_html = html_list[i]

                # 计算相似度
                similarity = compare_text(table_html, candidate_html, embeddings)

                if similarity > similarity_threshold:
                    # 找到匹配，获取对应的 DataFrame 内容
                    matched_df = df_list[i]
                    # 将 DataFrame 转换为更可读的文本格式
                    best_match_content = df_to_readable_text(matched_df)
                    html_pointer = i + 1  # 移动指针到下一位置
                    break

            if best_match_content:
                placeholders.append(f"<p>表{table_counter}：{best_match_content}</p>")
            else:
                # 如果没找到匹配，使用原始表格的简化文本
                table_text = extract_table_text_from_html(table)
                placeholders.append(f"<p>表{table_counter}：{table_text}</p>")

            table_counter += 1

        return "\n".join(placeholders)

    pattern = re.compile(r"<html>.*?</html>", re.DOTALL | re.IGNORECASE)
    new_content = pattern.sub(html_replacer, content)

    return new_content


def df_to_readable_text(df: pd.DataFrame) -> str:
    """
    将 DataFrame 转换为可读的文本格式，处理合并单元格情况
    """
    if df.empty:
        return "空表格"

    # 创建一个副本来避免修改原始数据
    df_copy = df.copy()

    # 处理垂直合并（向下填充）
    for col in df_copy.columns:
        df_copy[col] = df_copy[col].ffill()

    # 处理水平合并（向右填充）
    for idx in df_copy.index:
        row_values = df_copy.loc[idx].values
        filled_row = []
        last_valid_value = None

        for value in row_values:
            if pd.notna(value) and str(value).strip() != "":
                last_valid_value = value
                filled_row.append(value)
            else:
                # 如果当前值为空，用前一个有效值填充
                filled_row.append(last_valid_value if last_valid_value is not None else "空")

        df_copy.loc[idx] = filled_row

    # 转换为可读文本
    text_parts = []
    for index, row in df_copy.iterrows():
        row_text = " | ".join([str(cell) if pd.notna(cell) and str(cell).strip() != "" else "空" for cell in row])
        text_parts.append(row_text)

    return "\n".join(text_parts)


def extract_table_text_from_html(table_element) -> str:
    """
    从 HTML 表格元素中提取纯文本内容
    """
    rows = table_element.find_all("tr")
    text_parts = []

    for row in rows:
        cells = row.find_all(["td", "th"])
        cell_texts = [cell.get_text(strip=True) for cell in cells]
        if any(cell_texts):  # 只添加非空行
            text_parts.append(" | ".join(cell_texts))

    return "\n".join(text_parts) if text_parts else "空表格"


def compare_text(text1: str, text2: str, embeddings, max_length: int = 512) -> float:
    """
    计算两个文本的语义相似度，自动切断过长输入
    """
    try:
        text1_clean = strip_html_all(text1)
        text2_clean = strip_html_all(text2)

        # 如果文本过短，返回低相似度
        if len(text1_clean) < 10 or len(text2_clean) < 10:
            return 0.0

        # 安全切断，防止超过模型 token 限制
        text1_clean = text1_clean[:max_length]
        text2_clean = text2_clean[:max_length]

        # 生成向量
        vec1 = embeddings.embed_query(text1_clean)
        vec2 = embeddings.embed_query(text2_clean)

        # 转换为 numpy array
        vec1 = np.array(vec1).reshape(1, -1)
        vec2 = np.array(vec2).reshape(1, -1)

        # 计算余弦相似度
        similarity = cosine_similarity(vec1, vec2)[0][0]
        return similarity

    except Exception as e:
        logger.warning(f"相似度计算失败: {e}")
        return 0.0



def strip_html_all(text: str, method: str = "re") -> str:
    """
    去除文本中的 HTML 标签、换行符、制表符和多余空格

    :param text: 输入 HTML 文本
    :param method: 去标签方式，"re" 使用正则，"bs4" 使用 BeautifulSoup
    :return: 纯文本（无 HTML 标签、无换行符、无制表符、无多余空格）
    """
    if method == "bs4":
        clean = BeautifulSoup(text, "html.parser").get_text()
    elif method == "re":
        clean = re.sub(r"<[^>]+>", "", text)
    else:
        raise ValueError("method 必须是 're' 或 'bs4'")

    # 去掉换行符、制表符、多余空格
    clean = re.sub(r"\s+", "", clean)
    return clean

def _load_doc(doc_path: str):
    """
    转化为docx
    """
    logger.info("Converting doc file to docx file...")
    docx_path = doc_path + "x"
    subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--invisible",
            "--convert-to",
            "docx",
            "--outdir",
            os.path.dirname(docx_path),
            doc_path,
        ],
        check=True,
    )
    logger.info("Converted doc file to docx file.")
    text,doc = _load_docx(docx_path)
    os.remove(docx_path)
    return text,doc

def _load_docx(docx_path: str):
    """
    解析 docx，将表格替换为文字总结，段落按 run 拼接，保留多级序号和括号
    """
    doc = docx.Document(docx_path)
    text = ""
    table_count = 0

    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            # 按 run 拼接，保留所有内容
            para_text = "".join([run.text for run in block.runs]).strip()
            if para_text:  # 过滤空段落
                text += para_text + "\n"
        elif isinstance(block, Table):
            table_count += 1
            table_text = _extract_table_text(block)
            text += f"[表格{table_count}内容：{table_text}]\n"

    return text.strip(),doc

def _extract_table_text(table: Table) -> str:
    """
    提取表格内容为文本，兼容合并单元格，不保留样式，仅保留文字。
    输出格式示例：
        行1: 单元格1 | 单元格2 | 单元格3
        行2: ...
    """
    rows_text = []
    for row_idx, row in enumerate(table.rows, start=1):
        cell_texts = []
        for cell in row.cells:
            # 有时同一单元格对象在合并单元格中会重复引用，需要去重
            cell_content = " ".join(p.text.strip() for p in cell.paragraphs if p.text.strip())
            cell_texts.append(cell_content if cell_content else "（空）")
        # 去掉重复单元格（由 docx 合并引用导致）
        deduped_cells = []
        for i, c in enumerate(cell_texts):
            if i == 0 or c != cell_texts[i - 1]:
                deduped_cells.append(c)
        rows_text.append(" | ".join(deduped_cells))
    return "\n".join(rows_text)

async def get_file_extension(file_path: str) -> str:
    """
    获取文件扩展名（不带点），例如：
    """
    return os.path.splitext(file_path)[1].lstrip('.').lower()

async def save_upload_file(file, upload_dir: str = "./uploaded_files"):
    os.makedirs(upload_dir, exist_ok=True)

    # 原文件名（带后缀）
    original_name = file.filename
    base_name, ext = os.path.splitext(original_name)

    # 构造保存路径：例如 ./uploads/报告.docx
    save_path = os.path.join(upload_dir, f"{base_name}{ext}")

    # 若重名可加时间戳或随机后缀避免覆盖
    if os.path.exists(save_path):
        import time
        save_path = os.path.join(upload_dir, f"{base_name}_{int(time.time())}{ext}")

    # 写入文件
    contents = await file.read()
    with open(save_path, "wb") as f:
        f.write(contents)

    return save_path


def merge_duplicate_records(records):
    """
    合并阶段、步骤序号、作业点、作业事项任务相同的记录。
    其他字段内容累加（避免重复），多个项换行拼接并编号。
    """
    merged = {}
    unique_keys = ["阶段", "步骤序号", "作业点", "作业事项任务"]
    merge_fields = ["具体做什么", "做到什么程度", "特别风险", "特别风险管控", "所需材料物品等", "内容"]

    for rec in records:
        key = tuple(str(rec.get(k, "")).strip() if rec.get(k) is not None else "" for k in unique_keys)
        if key not in merged:
            merged[key] = {**rec}
            for f in merge_fields:
                merged[key][f] = [rec.get(f, "")] if rec.get(f) else []
        else:
            for f in merge_fields:
                val = rec.get(f)
                if val and val not in merged[key][f]:
                    merged[key][f].append(val)

    result = []
    for key, rec in merged.items():
        new_rec = {**rec}
        for f in merge_fields:
            vals = [str(v).strip() for v in rec[f] if v is not None and str(v).strip()]
            if not vals:
                new_rec[f] = ""
            elif len(vals) == 1:
                # 只有一个内容，不加编号
                new_rec[f] = vals[0]
            else:
                # 多个内容，加编号并换行拼接
                new_rec[f] = "\n".join([f"{i+1}. {v}" for i, v in enumerate(vals)])
        result.append(new_rec)

    return result

# --- 处理 JSON ---
def try_parse_json(text):
    """尝试直接 loads，失败则从 ```json ... ``` 里提取"""
    # 第一种情况：直接是 JSON
    try:
        return json.loads(text)
    except Exception:
        pass

    # 第二种情况：包在```json ... ```
    match = re.search(r"```json\s*(.*?)\s*```", text, re.S)
    if match:
        try:
            cleaned = match.group(1).strip()
            return json.loads(cleaned)
        except Exception:
            pass

    # 第三种情况：提取大括号内部的 JSON（兜底）
    match = re.search(r"(\{.*\})", text, re.S)
    if match:
        try:
            return json.loads(match.group(1))
        except Exception:
            pass

    raise ValueError(f"模型返回内容无法解析为 JSON：\n{text}")


async def detect_language_from_preview(preview_text, client: OpenAI):
    """
    根据预览内容识别语言
    """
    # 将预览列表转为字符串，取前1000个字符即可
    sample_content = str(preview_text)[:1000]
    prompt = PromptRegistry.get(PromptKey.GET_FILE_LANGUAGE_TYPE_235B,Lang.ZH,sample_content=sample_content)

    llm_config = get_dataprep_llm_config()
    extra_body = get_llm_extra_body(llm_config.model)

    try:
        response = await asyncio.to_thread(
            client.chat.completions.create,
            model=llm_config.model,
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
            temperature=0,
            extra_body=extra_body
        )
        res = json.loads(response.choices[0].message.content)
        return res.get("language", "zh")
    except Exception as e:
        logger.error(f"语言检测失败: {e}")
        return "zh"

def extract_json_from_response(response_text: str) -> dict | list:
    """
    从响应文本中提取 JSON 内容。

    支持以下常见模型返回形式：
    - 纯 JSON
    - ```json fenced code block
    - 带前后解释文字的 JSON
    - 带 <think>...</think> 前缀的 JSON
    - 空白或无效响应时抛出 ValueError

    Args:
        response_text: 包含 JSON 的响应文本

    Returns:
        解析后的 JSON 对象或数组
    """
    if response_text is None:
        raise ValueError("Response text is None")

    cleaned_text = response_text.strip()
    if not cleaned_text:
        raise ValueError("Response text is empty")

    if "</think>" in cleaned_text:
        cleaned_text = cleaned_text.split("</think>", 1)[1].strip()
    if not cleaned_text:
        raise ValueError("Response text is empty after removing think content")

    candidates: list[str] = []

    if "```json" in cleaned_text:
        json_start = cleaned_text.find("```json") + 7
        json_end = cleaned_text.find("```", json_start)
        if json_end != -1:
            candidates.append(cleaned_text[json_start:json_end].strip())
    elif "```" in cleaned_text:
        json_start = cleaned_text.find("```") + 3
        json_end = cleaned_text.find("```", json_start)
        if json_end != -1:
            candidates.append(cleaned_text[json_start:json_end].strip())

    candidates.append(cleaned_text)

    decoder = json.JSONDecoder()
    for candidate in candidates:
        if not candidate:
            continue
        try:
            return json.loads(candidate)
        except json.JSONDecodeError:
            pass

        for index, char in enumerate(candidate):
            if char not in "[{":
                continue
            try:
                parsed, _ = decoder.raw_decode(candidate[index:])
                if isinstance(parsed, (dict, list)):
                    return parsed
            except json.JSONDecodeError:
                continue

    preview = cleaned_text[:300].replace("\n", "\\n")
    raise ValueError(f"Failed to extract JSON from response: {preview}")

def build_chapters(word_tree: list):
    """
    只以一级标题为章节单位，每个章节包含完整的带标题结构正文
    """
    chapters = []

    for node in word_tree:
        chapters.append({
            "title": node["total_title"],
            "full_text": node["content"]
        })

    return chapters

import re

def clean_prefix(text: str) -> str:
    """
    去除开头的：
    1.  1、  1)  （1）
    要点1： 要点 1 : 要点1
    """
    text = text.strip()

    # 去掉类似 1.  1、  1)  (1)  （1） 这种编号
    text = re.sub(r'^[（(]?\d+[）).、]?\s*', '', text)

    # 去掉类似 要点1：  要点 12 :  要点1  要点 1
    text = re.sub(r'^要点\s*\d+\s*(?:[：:]?\s*)', '', text)

    return text.strip()
